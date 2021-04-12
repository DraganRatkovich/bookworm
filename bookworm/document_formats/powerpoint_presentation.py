# coding: utf-8

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from io import StringIO
from itertools import chain
from functools import cached_property
from more_itertools import all_equal
from bookworm.i18n import LocaleInfo
from bookworm.utils import NEWLINE
from bookworm.structured_text import SemanticElementType
from bookworm.document_formats.base import (
    BaseDocument,
    BasePage,
    BookMetadata,
    Section,
    Pager,
    DocumentCapability as DC,
    TreeStackBuilder,
    DocumentError,
    DocumentEncryptedError,
)
from bookworm.logger import logger


log = logger.getChild(__name__)
PP_HEADING_TYPES = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.SUBTITLE,
    PP_PLACEHOLDER.VERTICAL_TITLE,
}
PP_PLACEHOLDER_SEMANTIC_ELEMENTS = {h: SemanticElementType.HEADING_1 for h in PP_HEADING_TYPES}
PP_PLACEHOLDER_SEMANTIC_ELEMENTS |= {
    PP_PLACEHOLDER.TABLE: SemanticElementType.TABLE
}
PP_SHAPE_SEMANTIC_ELEMENTS = {
    MSO_SHAPE_TYPE.TABLE: SemanticElementType.TABLE,
}


class PowerpointSlide(BasePage):
    """Represents a slide in a PowerPoint presentation.""" 

    def __init__(self, slide, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.slide = slide
        self.text_buffer = StringIO(newline=NEWLINE)
        self.semantic_elements = {}
        self.style_elements = {}
        self.extract_slide_text_and_semantic(slide)
        self.extract_notes_slide(slide)

    def extract_slide_text_and_semantic(self, slide):
        shapes = (shape for shape in slide.shapes if shape.has_text_frame)
        for shape in shapes:
            text, is_list = self._get_shape_text(shape)
            start_pos = self.text_buffer.tell()
            self.text_buffer.write(text)
            stop_pos = self.text_buffer.tell()
            if is_list:
                self.semantic_elements.setdefault(SemanticElementType.LIST, []).append((start_pos, stop_pos))
            if shape.is_placeholder and (selm := PP_PLACEHOLDER_SEMANTIC_ELEMENTS.get(shape.placeholder_format.type)):
                self.semantic_elements.setdefault(selm, []).append((start_pos, stop_pos))
            elif (selm := PP_SHAPE_SEMANTIC_ELEMENTS.get(shape.shape_type)):
                self.semantic_elements.setdefault(selm, []).append((start_pos, stop_pos))

    def extract_notes_slide(self, slide):
        if not slide.has_notes_slide or not slide.notes_slide.notes_text_frame.text.strip():
            return
        self.text_buffer.write(NEWLINE + "-" * 10 + NEWLINE)
        nh_start_pos = self.text_buffer.tell()
        self.text_buffer.write(_("Slide Notes") + NEWLINE)
        nh_stop_pos = self.text_buffer.tell()
        self.semantic_elements.setdefault(SemanticElementType.HEADING_1, []).append((nh_start_pos, nh_stop_pos))
        self.extract_slide_text_and_semantic(slide.notes_slide)

    def get_text(self):
        return self.text_buffer.getvalue()

    def get_style_info(self) -> dict:
        return self.style_elements

    def get_semantic_structure(self) -> dict:
        return self.semantic_elements

    def _get_shape_text(self, shape):
        text = shape.text_frame.text.replace("\v", "\n")
        paragraphs = shape.text_frame.paragraphs
        parag_levels = [p.level for p in paragraphs if p.level > 0]
        is_list = (len(paragraphs) == len(parag_levels)) and all_equal(l for l in parag_levels)
        return text.strip(), is_list



class PowerpointPresentation(BaseDocument):

    format = "pptx"
    # Translators: the name of a document file format
    name = _("PowerPoint Presentation")
    extensions = ("*.pptx",)
    capabilities = DC.TOC_TREE | DC.METADATA | DC.STRUCTURED_NAVIGATION | DC.TEXT_STYLE

    def __len__(self):
        return self.num_slides

    def get_page(self, index):
        return PowerpointSlide(self.slides[index], self, index)

    def read(self):
        self.pptx_filepath = self.get_file_system_path()
        self.presentation = pptx.Presentation(self.pptx_filepath)
        self.slides = self.presentation.slides
        self.num_slides = len(self.slides)

    def close(self):
        super().close()

    @cached_property
    def language(self):
        if (lang := self.presentation.core_properties.language):
            try:
                return LocaleInfo(lang)
            except ValueError:
                pass
        return super().language

    @cached_property
    def toc_tree(self):
        root = Section(
            document=self,
            title=self.metadata.title,
            pager=Pager(first=0, last=self.num_slides - 1),
            level=1,
        )
        stack = TreeStackBuilder(root)
        for (idx, slide) in enumerate(self.slides):
            section_title = _("Slide {number}").format(number=idx+1)
            if (slide_title := self._get_slide_title(slide)):
                section_title = f"{section_title}: {slide_title}"
            stack.push(Section(
                document=self,
                title=section_title,
                pager=Pager(first=idx, last=idx),
                level=2
            ))
        return root

    @cached_property
    def metadata(self):
        props = self.presentation.core_properties
        presentation_title = props.title.strip() or self.pptx_filepath.stem.strip()
        return BookMetadata(
            title=presentation_title,
            author=props.author or "",
            publication_year=props.created or ""
        )

    def _get_slide_title(self, slide):
        try:
            first_shape = slide.shapes[0]
            if first_shape.is_placeholder and first_shape.placeholder_format.type in PP_HEADING_TYPES:
                return first_shape.text
        except IndexError:
            return